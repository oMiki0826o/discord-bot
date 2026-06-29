"""
core/ai/file_parser/document_parser.py

修正（文件類解析）：
- 涵蓋 pdf / docx / xlsx / pptx 文字提取，統一輸出純文字
- 本模組維持同步介面；CPU 密集的執行緒池調度統一由 __init__.py 的
  parse() 透過 asyncio.to_thread 處理，所有 parser（含本模組）保持
  相同的同步呼叫慣例，避免雙重包裝
- 各格式皆為可選依賴（pypdf / python-docx / openpyxl / python-pptx），
  缺少對應套件時回傳明確 error，不中斷整體流程
- PDF 僅提取文字層；掃描型 PDF（無文字層）回傳提示，不做 OCR（OCR 屬未來規劃）
"""

from __future__ import annotations

import logging
from pathlib import Path

from core.ai.file_parser.models import ParsedFile
from core.ai.file_parser.summary_builder import truncate

logger = logging.getLogger("bot.file_parser.document")

# Excel 單一工作表最大列數，避免大型試算表耗盡記憶體
_MAX_XLSX_ROWS = 200


# ── 主要入口 ─────────────────────────────────────────────────────────────

def parse(path: Path, filename: str, size_bytes: int) -> ParsedFile:
    """解析文件類檔案，回傳 ParsedFile（同步函式，由上層統一排程至執行緒池）。"""
    ext = path.suffix.lower()
    try:
        return _dispatch(path, filename, ext, size_bytes)
    except Exception as e:
        logger.debug("[document_parser] error file=%s: %s", filename, e)
        return ParsedFile(
            filename=filename, extension=ext,
            category="document", size_bytes=size_bytes,
            error=str(e),
        )


# ── 格式分派 ─────────────────────────────────────────────────────────────

def _dispatch(path: Path, filename: str, ext: str, size_bytes: int) -> ParsedFile:
    if ext == ".pdf":
        return _parse_pdf(path, filename, size_bytes)
    if ext == ".docx":
        return _parse_docx(path, filename, size_bytes)
    if ext == ".doc":
        return ParsedFile(
            filename=filename, extension=ext,
            category="document", size_bytes=size_bytes,
            error="不支援舊版 .doc 格式，請轉換為 .docx",
        )
    if ext == ".xlsx":
        return _parse_xlsx(path, filename, size_bytes)
    if ext == ".xls":
        return ParsedFile(
            filename=filename, extension=ext,
            category="document", size_bytes=size_bytes,
            error="不支援舊版 .xls 格式，請轉換為 .xlsx",
        )
    if ext == ".pptx":
        return _parse_pptx(path, filename, size_bytes)
    if ext == ".ppt":
        return ParsedFile(
            filename=filename, extension=ext,
            category="document", size_bytes=size_bytes,
            error="不支援舊版 .ppt 格式，請轉換為 .pptx",
        )
    return ParsedFile(
        filename=filename, extension=ext,
        category="document", size_bytes=size_bytes,
        error=f"不支援的文件格式：{ext}",
    )


# ── PDF ──────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path, filename: str, size_bytes: int) -> ParsedFile:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ParsedFile(
            filename=filename, extension=".pdf",
            category="document", size_bytes=size_bytes,
            error="缺少 pypdf 套件，無法解析 PDF",
        )

    reader = PdfReader(str(path))
    pages  = reader.pages
    texts:  list[str] = []

    for page in pages:
        text = page.extract_text() or ""
        if text.strip():
            texts.append(text)

    full_text = "\n\n".join(texts)

    if not full_text.strip():
        # 沒有文字層 → 可能是掃描 PDF，OCR 屬未來規劃，先明確告知
        return ParsedFile(
            filename=filename, extension=".pdf",
            category="document", size_bytes=size_bytes,
            content=f"[此 PDF 共 {len(pages)} 頁，未偵測到可提取的文字層，可能為掃描檔]",
            error=None,
        )

    content, truncated = truncate(full_text)
    return ParsedFile(
        filename=filename, extension=".pdf",
        category="document", size_bytes=size_bytes,
        content=f"[共 {len(pages)} 頁]\n\n{content}",
        truncated=truncated,
    )


# ── DOCX ─────────────────────────────────────────────────────────────────

def _parse_docx(path: Path, filename: str, size_bytes: int) -> ParsedFile:
    try:
        import docx
    except ImportError:
        return ParsedFile(
            filename=filename, extension=".docx",
            category="document", size_bytes=size_bytes,
            error="缺少 python-docx 套件，無法解析 Word 文件",
        )

    doc = docx.Document(str(path))

    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]

    # 表格內容也納入（Word 報告常見表格資料）
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    full_text = "\n".join(parts)
    content, truncated = truncate(full_text)
    return ParsedFile(
        filename=filename, extension=".docx",
        category="document", size_bytes=size_bytes,
        content=content, truncated=truncated,
    )


# ── XLSX ─────────────────────────────────────────────────────────────────

def _parse_xlsx(path: Path, filename: str, size_bytes: int) -> ParsedFile:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ParsedFile(
            filename=filename, extension=".xlsx",
            category="document", size_bytes=size_bytes,
            error="缺少 openpyxl 套件，無法解析 Excel 檔案",
        )

    # read_only=True：避免整份工作表載入記憶體
    wb = load_workbook(str(path), read_only=True, data_only=True)

    sections:  list[str] = []
    truncated = False

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sections.append(f"--- 工作表：{sheet_name} ---")
        row_count = 0
        for row in ws.iter_rows(max_row=_MAX_XLSX_ROWS, values_only=True):
            row_count += 1
            cells = [str(c) if c is not None else "" for c in row]
            sections.append(" | ".join(cells))
        # read_only 模式 max_row 可能不可靠，故以實際讀取列數判斷是否截斷
        if row_count >= _MAX_XLSX_ROWS:
            truncated = True
            sections.append(f"[此工作表僅顯示前 {_MAX_XLSX_ROWS} 列]")

    wb.close()

    full_text = "\n".join(sections)
    content, extra_cut = truncate(full_text)
    return ParsedFile(
        filename=filename, extension=".xlsx",
        category="document", size_bytes=size_bytes,
        content=content, truncated=truncated or extra_cut,
    )


# ── PPTX ─────────────────────────────────────────────────────────────────

def _parse_pptx(path: Path, filename: str, size_bytes: int) -> ParsedFile:
    try:
        from pptx import Presentation
    except ImportError:
        return ParsedFile(
            filename=filename, extension=".pptx",
            category="document", size_bytes=size_bytes,
            error="缺少 python-pptx 套件，無法解析 PowerPoint 檔案",
        )

    prs   = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide_texts:
            parts.append(f"--- 第 {i} 頁 ---\n" + "\n".join(slide_texts))

    full_text = "\n\n".join(parts)
    content, truncated = truncate(full_text)
    return ParsedFile(
        filename=filename, extension=".pptx",
        category="document", size_bytes=size_bytes,
        content=f"[共 {len(prs.slides)} 頁]\n\n{content}",
        truncated=truncated,
    )
